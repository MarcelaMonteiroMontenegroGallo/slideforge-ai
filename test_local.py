"""
Teste local do novo fluxo: template PPTX + preenchimento por IA.
Cria um template fictício da DataVision e preenche com conteúdo de proposta.
Roda: python test_local.py
"""
import sys, os, io
sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'app'))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from copy import deepcopy
from lxml import etree
from pptx.oxml.ns import qn

# ─── Conteúdo fictício (normalmente viria do Bedrock) ───
SLIDES_CONTENT = [
    {"type":"cover",   "title":"Plataforma de Analytics Unificada", "subtitle":"Proposta Comercial — TechRetail S.A."},
    {"type":"section", "title":"O Desafio", "subtitle":"Contexto e Problema"},
    {"type":"bullets", "title":"Situação Atual", "bullets":["Dados fragmentados em 4 sistemas sem integração","Time de BI leva 5 dias para responder perguntas","Falta de visibilidade em tempo real","Alta dependência de planilhas Excel"]},
    {"type":"section", "title":"Nossa Solução", "subtitle":"Databricks Lakehouse"},
    {"type":"bullets", "title":"Arquitetura Proposta", "bullets":["Data Lake centralizado com modelo Medallion","Integração automática com sistemas legados","Dashboard executivo com KPIs em tempo real","Genie AI para análise em linguagem natural"]},
    {"type":"metrics", "title":"Resultados Esperados", "metrics":[{"label":"Redução no tempo de análise","value":"80%"},{"label":"Redução de tickets de BI","value":"40%"},{"label":"ROI estimado","value":"18 meses"},{"label":"Lojas integradas","value":"280"}]},
    {"type":"bullets", "title":"Investimento", "bullets":["Fase 1 — Fundação (3 meses): R$ 180.000","Fase 2 — Analytics (2 meses): R$ 120.000","Fase 3 — IA e Genie (2 meses): R$ 80.000","Total: R$ 380.000 em 7 meses"]},
    {"type":"closing", "title":"Próximos Passos", "content":"Aprovação até 15/04 → Kick-off em 22/04 → Início Fase 1 em 01/05/2026\n\nVamos transformar seus dados em vantagem competitiva."},
]

PRIMARY = RGBColor(0x1a, 0x3a, 0x5c)
ACCENT  = RGBColor(0xff, 0x6b, 0x00)
WHITE   = RGBColor(255, 255, 255)
GRAY    = RGBColor(90, 90, 90)
LIGHT   = RGBColor(245, 247, 250)


def add_rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s


def add_placeholder_text(slide, text, l, t, w, h, size, color, bold=False, align=PP_ALIGN.LEFT, name=""):
    """Adiciona caixa de texto com nome de placeholder para o preenchimento automático."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.name = name  # nome usado para identificar o placeholder
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.color.rgb = color; run.font.bold = bold
    return tb


def create_template():
    """Cria um template PPTX fictício da DataVision com 8 slides em branco."""
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W, H  = Inches(13.33), Inches(7.5)

    layouts = [
        ("cover","capa"),("section","secao"),("bullets","conteudo"),
        ("section","secao"),("bullets","conteudo"),("metrics","metricas"),
        ("bullets","conteudo"),("closing","encerramento"),
    ]

    for slide_type, _ in layouts:
        sl = prs.slides.add_slide(blank)

        if slide_type == "cover":
            add_rect(sl, 0, 0, W, H, PRIMARY)
            add_rect(sl, 0, 0, Inches(0.5), H, ACCENT)
            add_rect(sl, Inches(0.8), Inches(3.3), Inches(9), Inches(0.05), ACCENT)
            # Logo placeholder (retângulo branco simulando logo)
            add_rect(sl, Inches(10.5), Inches(0.3), Inches(2.5), Inches(0.7), WHITE)
            add_placeholder_text(sl, "[TÍTULO DA APRESENTAÇÃO]", Inches(0.9), Inches(1.4), Inches(11), Inches(1.6), 36, WHITE, bold=True, name="title")
            add_placeholder_text(sl, "[Subtítulo / Cliente]",    Inches(0.9), Inches(3.5), Inches(10), Inches(0.8), 18, RGBColor(180,200,220), name="subtitle")
            add_placeholder_text(sl, "DataVision Consultoria",   Inches(0.9), Inches(5.7), Inches(7),  Inches(0.5), 14, ACCENT, bold=True)

        elif slide_type == "section":
            add_rect(sl, 0, 0, W, H, PRIMARY)
            add_rect(sl, 0, Inches(3.1), W, Inches(0.07), ACCENT)
            add_placeholder_text(sl, "[TÍTULO DA SEÇÃO]",  Inches(1), Inches(2.3), Inches(11), Inches(1.4), 40, WHITE, bold=True, align=PP_ALIGN.CENTER, name="title")
            add_placeholder_text(sl, "[Subtítulo]",        Inches(1), Inches(4.1), Inches(11), Inches(0.8), 18, RGBColor(170,190,210), align=PP_ALIGN.CENTER, name="subtitle")

        elif slide_type == "metrics":
            add_rect(sl, 0, 0, W, H, WHITE)
            add_rect(sl, 0, 0, W, Inches(1.15), PRIMARY)
            add_rect(sl, 0, Inches(1.15), Inches(3.5), Inches(0.06), ACCENT)
            add_placeholder_text(sl, "[TÍTULO]", Inches(0.4), Inches(0.18), Inches(12), Inches(0.8), 24, WHITE, bold=True, name="title")
            # 4 cards de métricas
            for j in range(4):
                x = Inches(0.5) + j * Inches(3.1)
                add_rect(sl, x, Inches(1.5), Inches(2.9), Inches(4.8), LIGHT)
                add_placeholder_text(sl, f"[VALOR {j+1}]", x+Inches(0.1), Inches(2.0), Inches(2.7), Inches(1.4), 40, ACCENT, bold=True, align=PP_ALIGN.CENTER, name=f"metric_value_{j}")
                add_placeholder_text(sl, f"[Label {j+1}]", x+Inches(0.1), Inches(4.0), Inches(2.7), Inches(1.0), 13, GRAY, align=PP_ALIGN.CENTER, name=f"metric_label_{j}")

        elif slide_type == "closing":
            add_rect(sl, 0, 0, W, H, PRIMARY)
            add_rect(sl, 0, 0, Inches(0.5), H, ACCENT)
            add_rect(sl, Inches(0.8), Inches(3.9), Inches(9), Inches(0.05), ACCENT)
            add_placeholder_text(sl, "[PRÓXIMOS PASSOS]", Inches(0.9), Inches(1.4), Inches(11), Inches(1.2), 32, WHITE, bold=True, name="title")
            add_placeholder_text(sl, "[Conteúdo de encerramento]", Inches(0.9), Inches(2.9), Inches(10), Inches(1.8), 17, RGBColor(190,210,225), name="body")
            add_placeholder_text(sl, "DataVision Consultoria", Inches(0.9), Inches(5.5), Inches(8), Inches(0.6), 15, ACCENT, bold=True)

        else:  # bullets / content
            add_rect(sl, 0, 0, W, H, WHITE)
            add_rect(sl, 0, 0, W, Inches(1.15), PRIMARY)
            add_rect(sl, 0, Inches(1.15), Inches(3.5), Inches(0.06), ACCENT)
            add_placeholder_text(sl, "[TÍTULO DO SLIDE]", Inches(0.4), Inches(0.18), Inches(12), Inches(0.8), 24, WHITE, bold=True, name="title")
            add_placeholder_text(sl, "• [Ponto 1]\n• [Ponto 2]\n• [Ponto 3]\n• [Ponto 4]", Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5), 16, GRAY, name="body")

        # Rodapé em todos os slides
        add_rect(sl, 0, Inches(7.2), W, Inches(0.3), LIGHT)
        add_placeholder_text(sl, "DataVision Consultoria  |  Confidencial", Inches(0.3), Inches(7.22), Inches(8), Inches(0.25), 9, GRAY)

    return prs


def fill_template(prs, slides_content):
    """Preenche os placeholders do template com o conteúdo."""
    for i, slide in enumerate(prs.slides):
        if i >= len(slides_content):
            break
        c = slides_content[i]
        title   = c.get("title", "")
        subtitle = c.get("subtitle", "")
        bullets = c.get("bullets", [])
        metrics = c.get("metrics", [])
        content = c.get("content", "")

        body_text = ""
        if bullets:
            body_text = "\n".join(f"• {b}" for b in bullets)
        elif content:
            body_text = content

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            n = shape.name.lower()

            if n == "title" and title:
                tf = shape.text_frame
                if tf.paragraphs and tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = title
                    # Remove parágrafos extras
                    while len(tf.paragraphs) > 1:
                        tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)

            elif n == "subtitle" and subtitle:
                tf = shape.text_frame
                if tf.paragraphs and tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = subtitle

            elif n == "body" and body_text:
                tf = shape.text_frame
                lines = body_text.split("\n")
                # Preenche primeiro parágrafo
                if tf.paragraphs and tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = lines[0]
                # Remove parágrafos extras e adiciona novos
                while len(tf.paragraphs) > 1:
                    tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
                for line in lines[1:]:
                    new_p = deepcopy(tf.paragraphs[0]._p)
                    tf.paragraphs[0]._p.addnext(new_p)
                    last = tf.paragraphs[-1]
                    if last.runs:
                        last.runs[0].text = line

            elif n.startswith("metric_value_") and metrics:
                idx = int(n.split("_")[-1])
                if idx < len(metrics):
                    tf = shape.text_frame
                    val = f"{metrics[idx].get('value','')}"
                    if tf.paragraphs and tf.paragraphs[0].runs:
                        tf.paragraphs[0].runs[0].text = val

            elif n.startswith("metric_label_") and metrics:
                idx = int(n.split("_")[-1])
                if idx < len(metrics):
                    tf = shape.text_frame
                    lbl = metrics[idx].get("label", "")
                    if tf.paragraphs and tf.paragraphs[0].runs:
                        tf.paragraphs[0].runs[0].text = lbl

    return prs


if __name__ == "__main__":
    print("1. Criando template DataVision...")
    prs_template = create_template()
    prs_template.save("datavision-TEMPLATE.pptx")
    print("   Salvo: datavision-TEMPLATE.pptx (template em branco com identidade visual)")

    print("2. Preenchendo template com conteúdo da proposta...")
    prs_filled = fill_template(prs_template, SLIDES_CONTENT)
    prs_filled.save("datavision-PROPOSTA-PREENCHIDA.pptx")
    print("   Salvo: datavision-PROPOSTA-PREENCHIDA.pptx")

    print(f"\nSlides: {len(SLIDES_CONTENT)}")
    print("Abra os dois arquivos para comparar:")
    print("  - datavision-TEMPLATE.pptx     → template em branco com identidade visual")
    print("  - datavision-PROPOSTA-PREENCHIDA.pptx → preenchido com conteúdo da proposta")
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'app'))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io

# ─── Cores da empresa fictícia DataVision ───
PRIMARY_HEX = "#1a3a5c"
ACCENT_HEX  = "#ff6b00"
COMPANY     = "DataVision Consultoria"
CLIENT      = "TechRetail S.A."

# ─── Conteúdo fictício (normalmente viria do Bedrock) ───
SLIDES = [
    {
        "type": "cover",
        "title": "Plataforma de Analytics Unificada",
        "subtitle": "Proposta Comercial — Transformação Digital de Dados",
    },
    {
        "type": "section",
        "title": "O Desafio",
        "subtitle": "Contexto e Problema",
    },
    {
        "type": "content",
        "title": "Situação Atual",
        "bullets": [
            "Dados fragmentados em 4 sistemas diferentes sem integração",
            "Time de BI leva 5 dias para responder perguntas estratégicas",
            "Falta de visibilidade em tempo real impacta decisões de estoque",
            "Campanhas promocionais baseadas em dados desatualizados",
            "Alta dependência de planilhas Excel para análises críticas",
        ],
    },
    {
        "type": "section",
        "title": "Nossa Solução",
        "subtitle": "Databricks Lakehouse + Analytics em Tempo Real",
    },
    {
        "type": "content",
        "title": "Arquitetura Proposta",
        "bullets": [
            "Data Lake centralizado com modelo Medallion (Bronze/Silver/Gold)",
            "Integração automática com os 4 sistemas legados via pipelines",
            "Dashboard executivo com KPIs atualizados em tempo real",
            "Genie AI para análise em linguagem natural pelo celular",
            "Governança centralizada com Unity Catalog e controle de acesso",
        ],
    },
    {
        "type": "metrics",
        "title": "Resultados Esperados",
        "metrics": [
            {"label": "Redução no tempo de análise", "value": "80%", "unit": ""},
            {"label": "Redução de tickets de BI",    "value": "40%", "unit": ""},
            {"label": "ROI estimado",                "value": "18",  "unit": "meses"},
            {"label": "Lojas integradas",            "value": "280", "unit": "lojas"},
        ],
    },
    {
        "type": "content",
        "title": "Fases de Implementação",
        "bullets": [
            "Fase 1 — Fundação (3 meses): Data Lake, pipelines e governança — R$ 180.000",
            "Fase 2 — Analytics (2 meses): Dashboards e KPIs em tempo real — R$ 120.000",
            "Fase 3 — IA e Genie (2 meses): Análise em linguagem natural — R$ 80.000",
            "Total: R$ 380.000 em 7 meses com equipe dedicada",
        ],
    },
    {
        "type": "content",
        "title": "Por que a DataVision?",
        "bullets": [
            "8 anos de experiência em arquitetura de dados moderna",
            "Mais de 60 projetos entregues no Brasil",
            "Parceira oficial Databricks e AWS",
            "Equipe certificada: 9x AWS + 2x Databricks por arquiteto",
            "Metodologia ágil com entregas quinzenais e transparência total",
        ],
    },
    {
        "type": "closing",
        "title": "Próximos Passos",
        "content": "Aprovação até 15/04 → Kick-off em 22/04 → Início Fase 1 em 01/05/2026\n\nVamos transformar seus dados em vantagem competitiva.",
    },
]


def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))


def add_rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def add_text(slide, text, l, t, w, h, size, color, bold=False, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size   = Pt(size)
    run.font.color.rgb = color
    run.font.bold   = bold
    run.font.italic = italic
    return tb


def build_ppt():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    primary = hex_to_rgb(PRIMARY_HEX)
    accent  = hex_to_rgb(ACCENT_HEX)
    white   = RGBColor(255, 255, 255)
    gray    = RGBColor(90, 90, 90)
    light   = RGBColor(245, 247, 250)
    dark    = RGBColor(30, 30, 30)

    blank = prs.slide_layouts[6]
    total = len(SLIDES)

    for i, data in enumerate(SLIDES):
        t    = data.get("type", "content")
        sl   = prs.slides.add_slide(blank)
        W, H = Inches(13.33), Inches(7.5)

        if t == "cover":
            add_rect(sl, 0, 0, W, H, primary)
            add_rect(sl, 0, 0, Inches(0.5), H, accent)
            add_rect(sl, Inches(0.8), Inches(3.3), Inches(9), Inches(0.05), accent)
            add_text(sl, data["title"],    Inches(0.9), Inches(1.4), Inches(11), Inches(1.6), 40, white, bold=True)
            add_text(sl, data.get("subtitle",""), Inches(0.9), Inches(3.5), Inches(10), Inches(0.8), 18, RGBColor(180,200,220))
            add_text(sl, f"Para: {CLIENT}", Inches(0.9), Inches(5.1), Inches(7), Inches(0.5), 13, RGBColor(160,180,200))
            add_text(sl, COMPANY,           Inches(0.9), Inches(5.7), Inches(7), Inches(0.5), 14, accent, bold=True)

        elif t == "section":
            add_rect(sl, 0, 0, W, H, primary)
            add_rect(sl, 0, Inches(3.1), W, Inches(0.07), accent)
            add_text(sl, data["title"],    Inches(1), Inches(2.3), Inches(11), Inches(1.4), 44, white, bold=True, align=PP_ALIGN.CENTER)
            if data.get("subtitle"):
                add_text(sl, data["subtitle"], Inches(1), Inches(4.1), Inches(11), Inches(0.8), 18, RGBColor(170,190,210), align=PP_ALIGN.CENTER)

        elif t == "metrics":
            add_rect(sl, 0, 0, W, H, white)
            add_rect(sl, 0, 0, W, Inches(1.15), primary)
            add_rect(sl, 0, Inches(1.15), Inches(3.5), Inches(0.06), accent)
            add_text(sl, data["title"], Inches(0.4), Inches(0.18), Inches(12), Inches(0.8), 24, white, bold=True)

            metrics = data.get("metrics", [])
            cols    = min(len(metrics), 4)
            cw      = Inches(12) / cols
            for j, m in enumerate(metrics[:4]):
                x = Inches(0.5) + j * cw
                add_rect(sl, x, Inches(1.5), cw - Inches(0.2), Inches(4.8), light)
                add_text(sl, m.get("value",""),  x+Inches(0.1), Inches(2.0), cw-Inches(0.4), Inches(1.4), 46, accent, bold=True, align=PP_ALIGN.CENTER)
                if m.get("unit"):
                    add_text(sl, m["unit"], x+Inches(0.1), Inches(3.3), cw-Inches(0.4), Inches(0.5), 14, gray, align=PP_ALIGN.CENTER)
                add_text(sl, m.get("label",""), x+Inches(0.1), Inches(4.0), cw-Inches(0.4), Inches(1.0), 13, dark, align=PP_ALIGN.CENTER)

        elif t == "closing":
            add_rect(sl, 0, 0, W, H, primary)
            add_rect(sl, 0, 0, Inches(0.5), H, accent)
            add_rect(sl, Inches(0.8), Inches(3.9), Inches(9), Inches(0.05), accent)
            add_text(sl, data["title"],   Inches(0.9), Inches(1.4), Inches(11), Inches(1.2), 36, white, bold=True)
            add_text(sl, data.get("content",""), Inches(0.9), Inches(2.9), Inches(10), Inches(1.8), 17, RGBColor(190,210,225))
            add_text(sl, COMPANY, Inches(0.9), Inches(5.5), Inches(8), Inches(0.6), 15, accent, bold=True)

        else:  # content / bullets
            add_rect(sl, 0, 0, W, H, white)
            add_rect(sl, 0, 0, W, Inches(1.15), primary)
            add_rect(sl, 0, Inches(1.15), Inches(3.5), Inches(0.06), accent)
            add_text(sl, data["title"], Inches(0.4), Inches(0.18), Inches(12), Inches(0.8), 24, white, bold=True)

            bullets = data.get("bullets", [])
            content = data.get("content", "")
            if bullets:
                for j, b in enumerate(bullets[:6]):
                    y = Inches(1.4) + j * Inches(0.88)
                    add_rect(sl, Inches(0.5), y + Inches(0.18), Inches(0.13), Inches(0.13), accent)
                    add_text(sl, b, Inches(0.8), y, Inches(11.8), Inches(0.8), 16, gray)
            elif content:
                add_text(sl, content, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5), 16, gray)

        # Número de slide (exceto capa)
        if t != "cover":
            add_text(sl, f"{i+1} / {total}", Inches(12.0), Inches(7.05), Inches(1.2), Inches(0.4), 10, RGBColor(160,160,160), align=PP_ALIGN.RIGHT)

        # Rodapé
        add_text(sl, COMPANY, Inches(0.3), Inches(7.1), Inches(5), Inches(0.35), 9, RGBColor(150,150,150))

    return prs


if __name__ == "__main__":
    print("Gerando PPT de demonstração...")
    prs = build_ppt()
    out = "datavision-proposta-demo.pptx"
    prs.save(out)
    print(f"PPT salvo: {out}")
    print(f"Slides: {len(SLIDES)}")
    print(f"Empresa: {COMPANY}")
    print(f"Cores: primária={PRIMARY_HEX}, destaque={ACCENT_HEX}")
