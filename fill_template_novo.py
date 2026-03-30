"""
Preenche o datavision-TEMPLATE NOVO.pptx com conteúdo da proposta.
Mapeamento exato baseado na inspeção do template.
Roda: python ppt-generator/fill_template_novo.py
"""
from pptx import Presentation
from copy import deepcopy

TEMPLATE_FILE = "datavision-TEMPLATE NOVO.pptx"
OUTPUT_FILE   = "datavision-PROPOSTA-NOVO-PREENCHIDA.pptx"

# ─── Conteúdo da proposta (normalmente viria do Bedrock) ───
PROPOSTA = {
    "slide1": {
        "titulo":   "Plataforma de Analytics Unificada na AWS",
        "subtitulo": "Proposta Comercial — TechRetail S.A.",
        "data":     "Março 2026",
    },
    "slide2": {
        "subtitulo": "Visão geral do cenário atual",
    },
    "slide3": {
        "titulo":  "Situação Atual",
        "ponto1":  "Dados fragmentados em 4 sistemas sem integração centralizada",
        "ponto2":  "Time de BI leva em média 5 dias para responder perguntas estratégicas",
        "ponto3":  "Falta de visibilidade em tempo real impacta decisões de estoque",
        "ponto4":  "Alta dependência de planilhas Excel para análises críticas do negócio",
    },
    "slide4": {
        "subtitulo": "Amazon S3 + AWS Glue + Athena + QuickSight + Bedrock",
    },
    "slide5": {
        "titulo":  "Arquitetura Proposta na AWS",
        "ponto1":  "Data Lake centralizado no Amazon S3 com modelo Medallion (Bronze/Silver/Gold)",
        "ponto2":  "Integração com sistemas legados via AWS Glue ETL e AWS DMS",
        "ponto3":  "Dashboard executivo com Amazon QuickSight e KPIs via Amazon Athena",
        "ponto4":  "Amazon Bedrock para análise em linguagem natural e insights automáticos",
    },
    "slide6": {
        "valor1": "80%",   "label1": "Redução no tempo de análise",
        "valor2": "40%",   "label2": "Redução de tickets de BI",
        "valor3": "18m",   "label3": "ROI estimado",
        "valor4": "280",   "label4": "Lojas integradas",
    },
    "slide7": {
        "fase1":        "Fundação — S3 Data Lake e AWS Glue",
        "prazo_fase1":  "3 meses",
        "valor_fase1":  "R$ 180.000",
        "fase2":        "Analytics — Athena e QuickSight",
        "prazo_fase2":  "2 meses",
        "valor_fase2":  "R$ 120.000",
        "fase3":        "IA — Amazon Bedrock e Insights",
        "prazo_fase3":  "2 meses",
        "valor_fase3":  "R$ 80.000",
        "valor_total":  "R$ 380.000",
        "prazo_total":  "7 meses",
    },
    "slide8": {
        "passo1": "Aprovação da proposta até 15/04/2026",
        "passo2": "Kick-off com equipe técnica em 22/04/2026",
        "passo3": "Início da Fase 1 em 01/05/2026",
    },
}


def set_text(shape, new_text):
    """Substitui o texto de uma shape preservando a formatação original."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    para = tf.paragraphs[0]
    if not para.runs:
        return
    # Preserva formatação do primeiro run e substitui o texto
    para.runs[0].text = new_text
    # Remove runs extras se houver
    for run in para.runs[1:]:
        run.text = ""


def fill_template(template_path, output_path, proposta):
    prs = Presentation(template_path)

    # ── SLIDE 1: CAPA ──────────────────────────────────────
    s1 = prs.slides[0]
    for shape in s1.shapes:
        if not shape.has_text_frame:
            continue
        n = shape.name
        t = shape.text_frame.paragraphs[0].text if shape.text_frame.paragraphs else ""

        if "[TÍTULO DA APRESENTAÇÃO]" in t:
            set_text(shape, proposta["slide1"]["titulo"])
        elif "[SUBTÍTULO / CLIENTE]" in t:
            set_text(shape, proposta["slide1"]["subtitulo"])
        elif "[DATA]" in t:
            set_text(shape, proposta["slide1"]["data"])

    # ── SLIDE 2: SEÇÃO CONTEXTO ────────────────────────────
    s2 = prs.slides[1]
    for shape in s2.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.paragraphs[0].text if shape.text_frame.paragraphs else ""
        if "[SUBTÍTULO DA SEÇÃO]" in t:
            set_text(shape, proposta["slide2"]["subtitulo"])

    # ── SLIDE 3: CONTEÚDO CONTEXTO (4 pontos) ─────────────
    s3 = prs.slides[2]
    pontos_contexto = [
        proposta["slide3"]["ponto1"],
        proposta["slide3"]["ponto2"],
        proposta["slide3"]["ponto3"],
        proposta["slide3"]["ponto4"],
    ]
    ponto_idx = 0
    for shape in s3.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.paragraphs[0].text if shape.text_frame.paragraphs else ""
        if "[TÍTULO DO SLIDE]" in t:
            set_text(shape, proposta["slide3"]["titulo"])
        elif "[CONTEXTO PONTO" in t and ponto_idx < len(pontos_contexto):
            set_text(shape, pontos_contexto[ponto_idx])
            ponto_idx += 1

    # ── SLIDE 4: SEÇÃO SOLUÇÃO ─────────────────────────────
    s4 = prs.slides[3]
    for shape in s4.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.paragraphs[0].text if shape.text_frame.paragraphs else ""
        if "[SUBTÍTULO DA SEÇÃO]" in t:
            set_text(shape, proposta["slide4"]["subtitulo"])

    # ── SLIDE 5: CONTEÚDO SOLUÇÃO (4 pontos) ──────────────
    s5 = prs.slides[4]
    pontos_solucao = [
        proposta["slide5"]["ponto1"],
        proposta["slide5"]["ponto2"],
        proposta["slide5"]["ponto3"],
        proposta["slide5"]["ponto4"],
    ]
    ponto_idx = 0
    for shape in s5.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.paragraphs[0].text if shape.text_frame.paragraphs else ""
        if "[TÍTULO DO SLIDE]" in t:
            set_text(shape, proposta["slide5"]["titulo"])
        elif "[SOLUÇÃO PONTO" in t and ponto_idx < len(pontos_solucao):
            set_text(shape, pontos_solucao[ponto_idx])
            ponto_idx += 1

    # ── SLIDE 6: MÉTRICAS (4 valores + labels) ─────────────
    s6 = prs.slides[5]
    metricas = {
        "[VALOR 1]":  proposta["slide6"]["valor1"],
        "[LABEL 1]":  proposta["slide6"]["label1"],
        "[VALOR 2]":  proposta["slide6"]["valor2"],
        "[LABEL 2]":  proposta["slide6"]["label2"],
        "[VALOR 3]":  proposta["slide6"]["valor3"],
        "[LABEL 3]":  proposta["slide6"]["label3"],
        "[VALOR 4]":  proposta["slide6"]["valor4"],
        "[LABEL 4]":  proposta["slide6"]["label4"],
    }
    for shape in s6.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.paragraphs[0].text if shape.text_frame.paragraphs else ""
        for placeholder, value in metricas.items():
            if placeholder in t:
                set_text(shape, value)
                break

    # ── SLIDE 7: INVESTIMENTO (3 fases + total) ────────────
    s7 = prs.slides[6]
    inv = proposta["slide7"]
    inv_map = {
        "[FASE 1]":        inv["fase1"],
        "[PRAZO FASE 1]":  inv["prazo_fase1"],
        "[VALOR FASE 1]":  inv["valor_fase1"],
        "[FASE 2]":        inv["fase2"],
        "[PRAZO FASE 2]":  inv["prazo_fase2"],
        "[VALOR FASE 2]":  inv["valor_fase2"],
        "[FASE 3]":        inv["fase3"],
        "[PRAZO FASE 3]":  inv["prazo_fase3"],
        "[VALOR FASE 3]":  inv["valor_fase3"],
    }
    for shape in s7.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.paragraphs[0].text if shape.text_frame.paragraphs else ""
        for placeholder, value in inv_map.items():
            if placeholder in t:
                set_text(shape, value)
                break
        # Total (texto composto)
        if "[VALOR TOTAL] em [PRAZO TOTAL]" in t:
            set_text(shape, f"{inv['valor_total']} em {inv['prazo_total']}")

    # ── SLIDE 8: PRÓXIMOS PASSOS ───────────────────────────
    s8 = prs.slides[7]
    passos = {
        "[PASSO 1]": proposta["slide8"]["passo1"],
        "[PASSO 2]": proposta["slide8"]["passo2"],
        "[PASSO 3]": proposta["slide8"]["passo3"],
    }
    for shape in s8.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.paragraphs[0].text if shape.text_frame.paragraphs else ""
        for placeholder, value in passos.items():
            if placeholder in t:
                set_text(shape, value)
                break

    prs.save(output_path)
    print(f"Salvo: {output_path}")


if __name__ == "__main__":
    print(f"Preenchendo template: {TEMPLATE_FILE}")
    fill_template(TEMPLATE_FILE, OUTPUT_FILE, PROPOSTA)
    print("Pronto! Abra o arquivo para ver o resultado.")
